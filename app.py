import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import google.generativeai as genai
import json
import io

# --- 1. CONFIGURATION & API SETUP ---
st.set_page_config(page_title="Weekly Math Lesson Plan Generator", page_icon="📐", layout="wide")

# Retrieve Gemini API key from Streamlit Secrets
if "GEMINI_API_KEY" in st.secrets:
    genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
else:
    st.error("Missing Gemini API Key! Please add 'GEMINI_API_KEY' to your Streamlit Settings > Secrets.")

# --- 2. POWERPOINT GENERATOR FUNCTION ---
def create_presentation(weekly_data, standard, topic):
    prs = Presentation()
    # Use standard 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    PRIMARY_COLOR = RGBColor(0x11, 0x18, 0x27)
    TEXT_COLOR = RGBColor(0x37, 0x41, 0x51)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    slide_types = [
        "Title", "Standard & Substandard", "Learning Target", "Opening (Warm-Up)", 
        "Work Session", "Formative Assessment", "GA Milestones Connection", 
        "Remediation", "Closing Activity", "Interdisciplinary Connection", "Project-Based Learning"
    ]
    
    for day in ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday"]:
        day_data = weekly_data.get(day, {})
        
        for slide_name in slide_types:
            slide = prs.slides.add_slide(blank_slide_layout)
            
            if slide_name != "Title":
                header_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(1.0))
                tf_header = header_box.text_frame
                tf_header.word_wrap = True
                p_day = tf_header.paragraphs[0]
                p_day.text = f"{day.upper()} • {topic}"
                p_day.font.size = Pt(12)
                p_day.font.bold = True
                p_day.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)
                
                p_title = tf_header.add_paragraph()
                p_title.text = slide_name
                p_title.font.size = Pt(28)
                p_title.font.bold = True
                p_title.font.color.rgb = PRIMARY_COLOR
                
                content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0))
                tf_content = content_box.text_frame
                tf_content.word_wrap = True
                
                content_text = day_data.get(slide_name, "No content generated.")
                lines = content_text.split('\n')
                for i, line in enumerate(lines):
                    if i == 0:
                        p = tf_content.paragraphs[0]
                    else:
                        p = tf_content.add_paragraph()
                    p.text = line
                    p.font.size = Pt(18)
                    p.font.color.rgb = TEXT_COLOR
                    if line.strip().startswith(('-', '*', '1.', '2.', '3.')):
                        p.level = 1
            else:
                title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.33), Inches(3.0))
                tf_title = title_box.text_frame
                tf_title.word_wrap = True
                
                p1 = tf_title.paragraphs[0]
                p1.text = f"{day}: {day_data.get('Title', topic)}"
                p1.font.size = Pt(44)
                p1.font.bold = True
                p1.font.color.rgb = PRIMARY_COLOR
                
                p2 = tf_title.add_paragraph()
                p2.text = f"Standard: {standard}\nTopic: {topic}"
                p2.font.size = Pt(20)
                p2.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
                
    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output

# --- 3. LLM PROMPT GENERATION ENGINE ---
def generate_day_plan(day, standard, substandard, topic):
    model = genai.GenerativeModel("gemini-pro")
    
    prompt = f"""
    You are an expert curriculum writer specializing in Georgia Mathematics standards.
    Generate a highly specific, actionable daily lesson plan for {day} based on the following details:
    - Standard: {standard}
    - Substandard: {substandard}
    - Specific Topic: {topic}
    
    You must output your response strictly as a JSON object with exactly these 11 keys. Do not include markdown formatting or backticks around the JSON. The values must contain explicit, robust educational content, sample problems, and strategies instead of generic outlines.
    
    {{
        "Title": "A descriptive, engaging lesson title for {day}",
        "Standard & Substandard": "Explicitly restate {standard} and {substandard} and how it applies to today's topic.",
        "Learning Target": "Write a clear, measurable 'I can' statement for the student.",
        "Opening (Warm-Up)": "A 5-10 minute diagnostic hook, warm-up question, or error analysis task.",
        "Work Session": "Detailed direct instruction concepts followed by a collaborative group task or independent practice activity.",
        "Formative Assessment": "A specific check for understanding strategy or 2 quick questions to evaluate learning mid-class.",
        "GA Milestones Connection": "A sample multiple-choice or multi-part problem designed exactly like a Georgia Milestones Assessment question, including an explanation of the correct answer.",
        "Remediation": "Specific modifications, scaffolding strategies, or concrete manipulatives/visual models to help struggling learners grasp this exact concept.",
        "Closing Activity": "A 5-minute wrap-up activity or a specific exit ticket question.",
        "Interdisciplinary Connection": "A concrete explanation of how this exact math topic connects to another field (e.g., a science application, historical context, or financial literacy concept).",
        "Project-Based Learning": "A description of a weekly mini-project milestone or hands-on application related to this topic that students can work on."
    }}
    """
    
    response = model.generate_content(prompt)
    try:
        clean_text = response.text.strip().strip("```json").strip("```").strip()
        return json.loads(clean_text)
    except Exception as e:
        return {k: f"Error generating content: {str(e)}. Raw text: {response.text}" for k in [
            "Title", "Standard & Substandard", "Learning Target", "Opening (Warm-Up)", 
            "Work Session", "Formative Assessment", "GA Milestones Connection", 
            "Remediation", "Closing Activity", "Interdisciplinary Connection", "Project-Based Learning"
        ]}

# --- 4. STREAMLIT USER INTERFACE ---
st.title("📐 Weekly Math Lesson Plan Presentation Generator")
st.write("Input your parameters below to instantly generate an entire 5-day, 55-slide aligned curriculum deck.")

# Nested Dictionary for Grade -> Standard -> Substandard
ga_standards = {
    "6th Grade": {
        "6.NR.1: Number System": ["6.NR.1.1", "6.NR.1.2"],
        "6.NR.2: Operations": ["6.NR.2.1", "6.NR.2.2"]
    },
    "7th Grade": {
        "7.NR.1: Rational Numbers": ["7.NR.1.1", "7.NR.1.2"],
        "7.PAR.2: Expressions": ["7.PAR.2.1", "7.PAR.2.2"]
    },
    "8th Grade": {
        "8.NR.1: Rational and Irrational Numbers": [
            "8.NR.1.1: Distinguish between rational and irrational numbers",
            "8.NR.1.2: Approximate irrational numbers to compare size"
        ],
        "8.NR.2: Integer Exponents and Scientific Notation": [
            "8.NR.2.1: Apply properties of integer exponents",
            "8.NR.2.2: Use square and cube root symbols",
            "8.NR.2.3: Estimate using scientific notation",
            "8.NR.2.4: Perform operations with scientific notation"
        ],
        "8.PAR.3: Linear Equations and Inequalities": [
            "8.PAR.3.1: Solve multi-step linear equations",
            "8.PAR.3.2: Solve multi-step linear inequalities"
        ],
        "8.PAR.4: Systems of Linear Equations": [
            "8.PAR.4.1: Understand solutions to systems of equations",
            "8.PAR.4.2: Solve systems of equations algebraically"
        ],
        "8.FGR.5: Functions": [
            "8.FGR.5.1: Understand a function is a rule",
            "8.FGR.5.2: Compare properties of two functions",
            "8.FGR.5.3: Interpret equation y = mx + b"
        ],
        "8.FGR.6: Transformations and Geometry": [
            "8.FGR.6.1: Verify properties of rotations, reflections, translations",
            "8.FGR.6.2: Understand congruence through transformations",
            "8.FGR.6.3: Describe effects of dilations"
        ],
        "8.FGR.7: Pythagorean Theorem": [
            "8.FGR.7.1: Explain a proof of the Pythagorean Theorem",
            "8.FGR.7.2: Apply theorem to find unknown side lengths",
            "8.FGR.7.3: Apply theorem to find distance on coordinate plane"
        ],
        "8.DSR.8: Bivariate Data": [
            "8.DSR.8.1: Construct and interpret scatter plots",
            "8.DSR.8.2: Use straight lines to model relationships",
            "8.DSR.8.3: Construct and interpret two-way tables"
        ]
    }
}

st.divider()

col1, col2, col3, col4 = st.columns(4)

with col1:
    grade_list = list(ga_standards.keys())
    default_grade_idx = grade_list.index("8th Grade") if "8th Grade" in grade_list else 0
    selected_grade = st.selectbox("Select Grade", grade_list, index=default_grade_idx)

with col2:
    standard_list = list(ga_standards[selected_grade].keys())
    standard = st.selectbox("Select Primary Standard", standard_list)

with col3:
    substandard_list = ga_standards[selected_grade][standard]
    substandard = st.selectbox("Select Substandard", substandard_list)

with col4:
    topic = st.text_input("Specific Topic", placeholder="e.g., Scientific Notation", key="unique_topic_input")

st.divider()

# --- 5. SUBMISSION BUTTON & LOGIC ---
if st.button("Generate Complete Widescreen Slide Deck", type="primary"):
    if not standard or not topic:
        st.warning("Please provide a specific topic to generate the lesson plan.")
    else:
        weekly_plan = {}
        
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        days = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday"]
        
        for i, day in enumerate(days):
            status_text.write(f"Drafting alignment and content for {day}...")
            weekly_plan[day] = generate_day_plan(day, standard, substandard, topic)
            progress_bar.progress((i + 1) / len(days))
            
        status_text.write("Compiling all 55 slides into widescreen PowerPoint format...")
        
        pptx_data = create_presentation(weekly_plan, standard, topic)
        
        status_text.success("Curriculum generation complete!")
        
        st.download_button(
            label="📥 Download PowerPoint Deck (.pptx)",
            data=pptx_data,
            file_name=f"{topic.replace(' ', '_')}_Weekly_Lesson_Plan.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
